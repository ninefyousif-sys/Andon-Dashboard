#!/usr/bin/env python3
"""
A-Shop Morning Meeting Dashboard — Daily Update Script
Run each morning at 05:50 Mon-Fri via Windows Task Scheduler.

Sources:
  Power BI Desktop  → FTT, DPV, W&G DPV, item details (DAX queries)
  Downtime Excel    → HOP + DT area downtime (same as Andon Dashboard)
  PPT 2026 RTM      → Safety, Part Quality, Bodies OOF, Scrap (tables only)
"""

import openpyxl, warnings, datetime, json, re, shutil, os, sys, subprocess
warnings.filterwarnings('ignore')

# ── CONFIG ─────────────────────────────────────────────────────────────────────
BASE_ONEDRIVE = r"C:\Users\NYOUSIF\OneDrive - Volvo Cars\A Shop Production SI and Supervisors - General"

HOP_SRC = os.path.join(BASE_ONEDRIVE, r"Hop Line Downtime\HOP New Downtime Breakdown.xlsm")
DT_SRC  = os.path.join(BASE_ONEDRIVE, r"Downtime Tracker Logv6a.xlsm")
PPT_DIR = os.path.join(BASE_ONEDRIVE, r"2026 RTM")

DASH        = r"C:\Users\NYOUSIF\Desktop\AShop_Dashboard\morning_meeting_dashboard.html"
DASH_MOBILE = r"C:\Users\NYOUSIF\Desktop\AShop_Dashboard\morning_meeting_mobile.html"
WORK        = r"C:\Users\NYOUSIF\Desktop\AShop_Dashboard"
LOG         = os.path.join(WORK, "mm_update_log.txt")

GITHUB_REPO    = WORK  # same folder — git repo is here
GITHUB_ENABLED = True  # pushes morning_meeting_mobile.html to GitHub Pages

WK12_MON = datetime.date(2026, 3, 16)

# ── BOK / BOL OPR TABLE CONFIG ─────────────────────────────────────────────────
# Run:  python update_morning_meeting.py --discover-tables
# with BIW SQD Axxos open in Power BI Desktop.
# The log will print all table names + measures.  Fill in the correct values below.
#
# ── CONFIRMED: 'OPR BOK Card' is a MEASURES TABLE (no data rows).
#    Use OPR_IS_MEASURES_TABLE = True and set OPR_BOK_MEASURE / OPR_BOL_MEASURE
#    to the DAX measure names shown in --discover-tables output.
#
# Example for measures-table setup:
#   OPR_TABLE          = 'OPR BOK Card'
#   OPR_IS_MEASURES_TABLE = True
#   OPR_BOK_MEASURE    = 'BOK OPR'   # measure name (from --discover-tables)
#   OPR_BOL_MEASURE    = 'BOL OPR'   # measure name (from --discover-tables)
#   # OPR_DATE_TABLE tries common names automatically — override if needed:
#   OPR_DATE_TABLE     = None        # e.g. 'Date' or 'Date Table' or 'Calendar'
#
# Example for regular data-table setup:
#   OPR_TABLE          = 'A Shop BOK BOL OPR'
#   OPR_IS_MEASURES_TABLE = False
#   OPR_COL_DATE       = 'Date'
#   OPR_BOK_MEASURE    = 'BOK OPR'   # column name
#   OPR_BOL_MEASURE    = 'BOL OPR'   # column name
#
OPR_TABLE             = 'OPR BOK Card'  # confirmed via --discover-tables
OPR_IS_MEASURES_TABLE = True            # confirmed — no data columns, only measures
OPR_BOK_MEASURE       = 'OPR BOK Card'  # confirmed via --discover-tables (2026-03-21)
OPR_BOL_MEASURE       = 'OPR BOL Card'  # confirmed via --discover-tables (2026-03-21)
OPR_DATE_TABLE        = 'Date Table'    # confirmed present in port 58016
OPR_COL_DATE          = 'Date'          # (used only when OPR_IS_MEASURES_TABLE=False)
# Legacy aliases kept for backwards compat
OPR_COL_BOK = OPR_BOK_MEASURE
OPR_COL_BOL = OPR_BOL_MEASURE

# ── SCRAP (Power BI) ────────────────────────────────────────────────────────────
# Shows as "Scrap / Car" ($X.XX) in the Body Shop KPI dashboard.
# Leave blank — query_scrap() auto-discovers the measure across all common names.
# After first successful run, the log will print a HINT with the exact names to
# set here so future runs skip the auto-discovery loop.
SCRAP_TABLE        = ''    # e.g. 'Scrap Card'   — auto-discovered if blank
SCRAP_MEASURE      = ''    # e.g. 'Scrap/Car'    — auto-discovered if blank
SCRAP_DATE_TABLE   = ''    # leave blank to reuse OPR_DATE_TABLE automatically

# ── TARGETS ────────────────────────────────────────────────────────────────────
TARGETS = {
    'bok_opr':    {'tgt': 98.7,  'dir': 'ge'},
    'bol_opr':    {'tgt': 95.0,  'dir': 'ge'},
    'ashop_ftt':  {'tgt': 98.0,  'dir': 'ge'},
    'ashop_dpv':  {'tgt': 0.10,  'dir': 'le'},
    'wg_dpv':     {'tgt': 1.20,  'dir': 'le'},
    'wd6_ftt':    {'tgt': 100.0, 'dir': 'ge'},
    'wd6_dpv':    {'tgt': 0.12,  'dir': 'le'},
    'cal_ftt':    {'tgt': 100.0, 'dir': 'ge'},
    'final1_ftt': {'tgt': 100.0, 'dir': 'ge'},
    'final2_ftt': {'tgt': 100.0, 'dir': 'ge'},
    'scrap_car':  {'tgt': 2.90,  'dir': 'le'},
}

# ── AREA MAPPING (Downtime) ────────────────────────────────────────────────────
HOP_AREA = {'285':'upperbody','286':'upperbody','287':'hangon','197':'hangon','198':'hangon','155':'bodysides'}
DT_AREA  = {'231':'underbody','232':'underbody','233':'underbody','234':'underbody',
            '235':'underbody','236':'underbody','237':'underbody','138':'underbody',
            '258':'bodysides','155':'bodysides'}
PLANNED_KW = {'pm','planned','planned maintenance','preventive','scheduled','maintenance pm'}

# ── HELPERS ────────────────────────────────────────────────────────────────────
def log(msg):
    ts = datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
    line = f"[{ts}] {msg}"
    print(line)
    with open(LOG, 'a', encoding='utf-8') as f: f.write(line+'\n')

def time_to_min(t):
    if t is None: return 0
    if isinstance(t, (datetime.datetime, datetime.time)): return t.hour*60+t.minute+t.second/60
    try:
        f = float(t); return f*1440 if 0 < f < 2 else f
    except: return 0

def dur_min(s, e):
    sm, em = time_to_min(s), time_to_min(e)
    return round(em-sm, 1) if em > sm else 0

def t_fmt(t):
    if isinstance(t, (datetime.datetime, datetime.time)): return f"{t.hour:02d}:{t.minute:02d}"
    return ''

def dt_str(t):
    if isinstance(t, datetime.datetime): return t.strftime('%m/%d/%Y %I:%M %p')
    if isinstance(t, datetime.date): return t.strftime('%m/%d/%Y')
    return str(t) if t else ''

def is_planned(resp, err):
    for v in [resp, err]:
        if v and any(k in str(v).lower() for k in PLANNED_KW): return True
    return False

def should_exclude(resp, err):
    r = str(resp).lower().strip() if resp else ''
    e = str(err).lower().strip()  if err  else ''
    if 'shop flow' in r or r in {'shop flow','shop flow '}: return True
    if 'blocked' in e or 'starved' in e: return True
    return False

def hop_code(area):
    a = str(area)
    for pre, c in [('285','285'),('286','286'),('287','287'),('198','198'),('155','155'),('197','197')]:
        if pre in a: return c
    if 'tailgate' in a.lower(): return '287'
    if 'roof'     in a.lower(): return '155'
    m = re.search(r'\b(\d{3})\b', a)
    return m.group(1) if m else None

def dt_code(area):
    a = str(area).lower()
    for pat, c in [('232','232'),('233','233'),('234','234'),('235','235'),('236','236'),
                   ('237','237'),('231','231'),('138','138'),('258','258'),('155','155'),('roof','155')]:
        if pat in a: return c
    m = re.search(r'\b(2\d{2})\b', str(area))
    return m.group(1) if m else None

def date_to_wk_day(d):
    if d.weekday() > 4: return None, None
    delta = (d - WK12_MON).days
    return f"WK{12 + delta//7}", d.weekday()+1

def prev_working_day(ref=None):
    d = (ref or datetime.date.today()) - datetime.timedelta(days=1)
    while d.weekday() > 4: d -= datetime.timedelta(days=1)
    return d


# ═══════════════════════════════════════════════════════════════════════════════
# POWER BI INTEGRATION
# Queries the open Power BI Desktop model via its local XMLA endpoint.
# Requires: pip install pyadomd  (wraps ADOMD.NET — ships with Power BI Desktop)
# ═══════════════════════════════════════════════════════════════════════════════

def find_pbi_ports():
    """Find ALL Analysis Services ports exposed by running Power BI Desktop instances.
    Returns a list of integer ports (may contain multiple if several PBI files are open)."""
    ports = []

    # Method 1: scan AnalysisServicesWorkspaces folder — PBI writes one port file per instance
    appdata = os.environ.get('LOCALAPPDATA', '')
    ws_root = os.path.join(appdata, 'Microsoft', 'Power BI Desktop',
                           'AnalysisServicesWorkspaces')
    if os.path.isdir(ws_root):
        for folder in sorted(os.listdir(ws_root)):
            port_file = os.path.join(ws_root, folder, 'Data', 'msmdsrv.port.txt')
            if os.path.exists(port_file):
                try:
                    port = int(open(port_file).read().strip())
                    if port not in ports:
                        ports.append(port)
                        log(f"Found PBI XMLA port: {port} (from {folder})")
                except: pass

    if ports:
        return ports

    # Method 2: PowerShell Get-WmiObject — collect ALL msmdsrv pids then ALL their ports
    try:
        ps_out = subprocess.check_output(
            ['powershell', '-NoProfile', '-Command',
             'Get-WmiObject Win32_Process | Where-Object {$_.Name -eq "msmdsrv.exe"} | Select-Object -ExpandProperty ProcessId'],
            text=True, timeout=10)
        pids = set(re.findall(r'\d{4,6}', ps_out))
        if pids:
            net_out = subprocess.check_output(['netstat', '-ano'], text=True, timeout=10)
            for line in net_out.splitlines():
                parts = line.split()
                if len(parts) >= 5 and parts[4] in pids and 'LISTENING' in line:
                    m = re.search(r':(\d{4,6})\s', line)
                    if m:
                        port = int(m.group(1))
                        if port not in ports:
                            ports.append(port)
                            log(f"Found PBI XMLA port via PowerShell/netstat: {port}")
    except: pass

    if ports:
        return ports

    # Method 3: fallback wmic scan (older Windows)
    try:
        out = subprocess.check_output(['netstat','-ano'], text=True, timeout=10)
        pids_raw = subprocess.check_output(
            ['wmic','process','where','name="msmdsrv.exe"','get','ProcessId'],
            text=True, timeout=10)
        pids = set(re.findall(r'\d{4,6}', pids_raw))
        for line in out.splitlines():
            parts = line.split()
            if len(parts) >= 5 and parts[4] in pids and 'LISTENING' in line:
                m = re.search(r':(\d{4,6})\s', line)
                if m:
                    port = int(m.group(1))
                    if port not in ports:
                        ports.append(port)
                        log(f"Found PBI XMLA port via wmic/netstat: {port}")
    except: pass

    if not ports:
        log("WARNING: Power BI Desktop not found / not running")
    return ports


def find_pbi_port():
    """Compatibility shim — returns first port found (use find_pbi_ports() for multi-instance)."""
    ports = find_pbi_ports()
    return ports[0] if ports else None


def run_dax(port, dax):
    """Run a DAX query against Power BI Desktop.
    Primary: PowerShell COM/ADODB (no Python packages needed — uses MSOLAP OLE DB
             provider installed with Power BI Desktop).
    Fallback: pyadomd if installed."""

    # ── Primary: PowerShell ADODB ─────────────────────────────────────────────
    tmp_dax = os.path.join(WORK, '_mm_dax_q.txt')
    tmp_out = os.path.join(WORK, '_mm_dax_r.json')
    tmp_ps  = os.path.join(WORK, '_mm_dax_run.ps1')

    ps_script = r"""param($Port, $DaxFile, $OutFile)
$ErrorActionPreference = 'Stop'
try {
    $dax  = Get-Content -Path $DaxFile -Raw -Encoding UTF8
    $conn = New-Object -ComObject ADODB.Connection
    $conn.Open("Provider=MSOLAP;Data Source=localhost:$Port;")
    $rs   = New-Object -ComObject ADODB.Recordset
    $rs.Open($dax, $conn)
    $fields = @()
    for ($i = 0; $i -lt $rs.Fields.Count; $i++) {
        $n = $rs.Fields.Item($i).Name
        if ($n -match '\[([^\]]+)\]$') { $n = $Matches[1] }
        $fields += $n
    }
    $rows = @()
    while (-not $rs.EOF) {
        $row = [ordered]@{}
        for ($i = 0; $i -lt $fields.Count; $i++) {
            $v = $rs.Fields.Item($i).Value
            if ($v -is [System.DBNull]) { $v = $null }
            $row[$fields[$i]] = $v
        }
        $rows += [PSCustomObject]$row
        $rs.MoveNext()
    }
    $rs.Close(); $conn.Close()
    if ($rows.Count -eq 0) { '[]' | Set-Content $OutFile -Encoding UTF8 }
    else { $rows | ConvertTo-Json -Depth 4 -Compress | Set-Content $OutFile -Encoding UTF8 }
} catch {
    Write-Host "DAX_PS_ERROR: $($_.Exception.Message)"
    '[]' | Set-Content $OutFile -Encoding UTF8
}
"""
    try:
        with open(tmp_dax, 'w', encoding='utf-8') as f: f.write(dax.strip())
        with open(tmp_ps,  'w', encoding='utf-8') as f: f.write(ps_script)
        subprocess.run(
            ['powershell', '-NoProfile', '-ExecutionPolicy', 'Bypass',
             '-File', tmp_ps, '-Port', str(port),
             '-DaxFile', tmp_dax, '-OutFile', tmp_out],
            capture_output=True, text=True, timeout=60)
        if os.path.exists(tmp_out):
            # utf-8-sig strips the BOM that Windows PowerShell 5.x adds when writing UTF-8
            with open(tmp_out, 'r', encoding='utf-8-sig') as f: raw = f.read().strip()
            if raw and raw != '[]':
                data = json.loads(raw)
                if isinstance(data, dict): data = [data]
                # clean column names: strip table prefix if present
                cleaned = []
                for row in data:
                    cleaned.append({re.sub(r'.*\[([^\]]+)\]$', r'\1', k): v
                                    for k, v in row.items()})
                return cleaned
    except Exception as e:
        log(f"DAX PowerShell error: {e}\nQuery: {dax[:120]}...")
    finally:
        for fp in [tmp_dax, tmp_out, tmp_ps]:
            try: os.remove(fp)
            except: pass

    # ── Fallback: pyadomd ─────────────────────────────────────────────────────
    try:
        import pyadomd
        conn_str = f"Provider=MSOLAP;Data Source=localhost:{port};"
        with pyadomd.Pyadomd(conn_str) as conn:
            with conn.cursor().execute(dax) as cur:
                cols = [c.name.split('[')[-1].rstrip(']') for c in cur.description]
                return [dict(zip(cols, row)) for row in cur.fetchall()]
    except ImportError:
        pass  # pyadomd not installed — PowerShell method already tried above
    except Exception as e:
        log(f"DAX pyadomd error: {e}\nQuery: {dax[:120]}...")

    return []


def fmt_dt(v):
    """Format a datetime value from DAX result to string.
    Handles: datetime objects, OData /Date(ms)/ strings, plain strings."""
    import datetime as _dt
    if v is None: return ''
    if hasattr(v, 'strftime'): return v.strftime('%m/%d/%Y %I:%M %p')
    s = str(v)
    # OData JSON date: /Date(1773947962000)/
    m = re.match(r'^/Date\((-?\d+)\)/$', s.strip())
    if m:
        try:
            ts = int(m.group(1)) / 1000.0
            return _dt.datetime.fromtimestamp(ts).strftime('%m/%d/%Y %I:%M %p')
        except Exception:
            pass
    return s


def discover_tables():
    """Print ALL table names (and columns for OPR/BOK/BOL tables) from every
    running Power BI Desktop instance.  Run once:
        python update_morning_meeting.py --discover-tables
    Then fill in OPR_TABLE / OPR_BOK_MEASURE / OPR_BOL_MEASURE at the top of this file."""
    ports = find_pbi_ports()
    if not ports:
        log("No PBI instances found.  Open BIW SQD Axxos in Power BI Desktop first.")
        return

    for port in ports:
        log(f"\n{'='*60}")
        log(f"=== Tables in PBI port {port} ===")
        log(f"{'='*60}")

        # DMV query — list all user tables in this model
        try:
            rows = run_dax(port,
                "SELECT [NAME],[DATABASE_NAME] FROM $SYSTEM.TMSCHEMA_TABLES "
                "WHERE [ISSHOWN] = TRUE()")
        except Exception as e:
            rows = []
            log(f"  DMV error: {e}")

        if not rows:
            # Fallback: try without filter
            try:
                rows = run_dax(port, "SELECT [NAME] FROM $SYSTEM.TMSCHEMA_TABLES")
            except Exception as e2:
                log(f"  DMV fallback error: {e2}")
                rows = []

        if not rows:
            log("  (no tables returned — model may not support DMV)")
            continue

        for r in rows:
            tname = r.get('NAME', r.get('name', ''))
            log(f"  TABLE: {tname}")

            # For any table with OPR / BOK / BOL / AVAIL in the name, also dump columns + measures
            if any(kw in tname.upper() for kw in ['OPR', 'BOK', 'BOL', 'AVAIL']):
                # Data columns
                try:
                    col_rows = run_dax(port,
                        f"SELECT [EXPLICIT_NAME],[DATATYPE_NAME] "
                        f"FROM $SYSTEM.TMSCHEMA_COLUMNS "
                        f"WHERE [TABLE_NAME] = '{tname}'")
                    for cr in col_rows:
                        cname = cr.get('EXPLICIT_NAME', cr.get('Name', ''))
                        ctype = cr.get('DATATYPE_NAME', cr.get('DataType', ''))
                        log(f"      COLUMN:  {cname}  ({ctype})")
                    if not col_rows:
                        log(f"      (no data columns — likely a measures-only table)")
                except Exception as ce:
                    log(f"      (could not list columns: {ce})")

                # Measures — use MDSCHEMA_MEASURES (uses MEASUREGROUP_NAME = table display name)
                meas_names = []
                try:
                    meas_rows = run_dax(port,
                        f"SELECT [MEASURE_NAME],[MEASUREGROUP_NAME] "
                        f"FROM $SYSTEM.MDSCHEMA_MEASURES "
                        f"WHERE [MEASUREGROUP_NAME] = '{tname}'")
                    meas_names = [mr.get('MEASURE_NAME', mr.get('measure_name', ''))
                                  for mr in meas_rows if mr.get('MEASURE_NAME') or mr.get('measure_name')]
                except Exception:
                    pass

                if not meas_names:
                    # Fallback: get all TMSCHEMA_MEASURES and cross-reference by table ID
                    try:
                        tbl_id_rows = run_dax(port,
                            f"SELECT [ID] FROM $SYSTEM.TMSCHEMA_TABLES WHERE [NAME] = '{tname}'")
                        tbl_id = (tbl_id_rows[0].get('ID') or tbl_id_rows[0].get('id')) if tbl_id_rows else None
                        if tbl_id is not None:
                            all_meas = run_dax(port,
                                "SELECT [NAME],[TABLEID] FROM $SYSTEM.TMSCHEMA_MEASURES")
                            meas_names = [r.get('NAME', r.get('name', ''))
                                          for r in all_meas
                                          if str(r.get('TABLEID', r.get('tableid', ''))) == str(tbl_id)]
                    except Exception:
                        pass

                for mname in meas_names:
                    log(f"      MEASURE: {mname}")
                if meas_names:
                    log(f"      >>> Set OPR_TABLE='{tname}', OPR_IS_MEASURES_TABLE=True")
                    log(f"          OPR_BOK_MEASURE='<BOK measure name above>'")
                    log(f"          OPR_BOL_MEASURE='<BOL measure name above>'")
                elif not col_rows:
                    log(f"      (no measures found via DMV — table may use hidden/implicit measures)")

        # ── Broad search: any measure with BOK / BOL / OPR in name (all tables) ──
        log(f"\n--- Searching ALL measures for BOK / BOL / OPR keywords ---")
        try:
            all_meas = run_dax(port,
                "SELECT [MEASURE_NAME],[MEASUREGROUP_NAME] FROM $SYSTEM.MDSCHEMA_MEASURES")
            found = [m for m in all_meas
                     if any(kw in str(m.get('MEASURE_NAME','') or '').upper()
                            for kw in ['BOK','BOL','OPR','AVAIL'])]
            for m in found:
                log(f"  FOUND: [{m.get('MEASURE_NAME','')}]  in table '{m.get('MEASUREGROUP_NAME','')}'")
            if not found:
                log("  (no BOK/BOL/OPR measures found by name search)")
        except Exception as be:
            log(f"  (broad measure search error: {be})")

        log(f"\n>>> ACTION: Fill in OPR_TABLE / OPR_BOK_MEASURE / OPR_BOL_MEASURE")
        log(f"            at the top of update_morning_meeting.py")
        log(f"            Set OPR_IS_MEASURES_TABLE=True if the table has MEASURE: lines above")


def query_opr(ports, report_date):
    """Query BOK/BOL OPR from any available PBI instance.

    Supports two modes controlled by OPR_IS_MEASURES_TABLE:

    True  (measures-container table like 'OPR BOK Card'):
        Uses CALCULATETABLE(ROW(...), DateTable[Date] = DATE(...))
        Tries several common date table names automatically.

    False (regular data table):
        Uses FILTER/SELECTCOLUMNS by date column.

    Returns {'bok_opr': float|None, 'bol_opr': float|None}, or None if not
    configured or no data found.
    """
    if not OPR_TABLE:
        log("  BOK/BOL OPR: OPR_TABLE not configured — run --discover-tables")
        return None

    bok_measure = OPR_BOK_MEASURE or OPR_COL_BOK
    bol_measure = OPR_BOL_MEASURE or OPR_COL_BOL
    if not bok_measure or not bol_measure:
        log("  BOK/BOL OPR: OPR_BOK_MEASURE / OPR_BOL_MEASURE not configured — run --discover-tables")
        return None

    yr, mo, dy = report_date.year, report_date.month, report_date.day

    def to_pct(v):
        if v is None: return None
        try:
            f = float(v)
            return round(f * 100, 2) if f <= 1.0 else round(f, 2)
        except (TypeError, ValueError):
            return None

    for port in ports:
        if OPR_IS_MEASURES_TABLE:
            # ── Measures-table path (e.g. 'OPR BOK Card') ─────────────────
            # CALCULATETABLE sets a date filter context so measures evaluate
            # against the correct day.  Try several common date table names.
            date_tables = []
            if OPR_DATE_TABLE:
                date_tables = [OPR_DATE_TABLE]
            else:
                date_tables = ['Date', 'Date Table', 'Calendar',
                               'DimDate', 'Dates', 'DateTable', 'Dim Date']

            for dt in date_tables:
                dax = (
                    f"EVALUATE CALCULATETABLE("
                    f"ROW("
                    f"\"bok_opr\", '{OPR_TABLE}'[{bok_measure}], "
                    f"\"bol_opr\", '{OPR_TABLE}'[{bol_measure}]),"
                    f"'{dt}'[Date] = DATE({yr},{mo},{dy}))"
                )
                try:
                    rows = run_dax(port, dax)
                    if rows:
                        r = rows[0]
                        bok = to_pct(r.get('bok_opr') or r.get('[bok_opr]'))
                        bol = to_pct(r.get('bol_opr') or r.get('[bol_opr]'))
                        if bok is not None or bol is not None:
                            log(f"  OPR (measures, date table='{dt}'): BOK={bok}%  BOL={bol}%")
                            return {'bok_opr': bok, 'bol_opr': bol}
                except Exception as e:
                    err_str = str(e)
                    if 'cannot be found' in err_str or 'not found' in err_str.lower():
                        continue   # wrong date-table name, try next
                    log(f"  OPR measures query error (port {port}, dt='{dt}'): {e}")

            # Last-resort: evaluate measure without date filter (returns overall value)
            dax_nf = (
                f"EVALUATE ROW("
                f"\"bok_opr\", '{OPR_TABLE}'[{bok_measure}], "
                f"\"bol_opr\", '{OPR_TABLE}'[{bol_measure}])"
            )
            try:
                rows = run_dax(port, dax_nf)
                if rows:
                    r = rows[0]
                    bok = to_pct(r.get('bok_opr') or r.get('[bok_opr]'))
                    bol = to_pct(r.get('bol_opr') or r.get('[bol_opr]'))
                    log(f"  OPR (measures, no date filter — overall value): BOK={bok}%  BOL={bol}%")
                    log(f"  WARNING: date filter failed — set OPR_DATE_TABLE to the correct date table name")
                    return {'bok_opr': bok, 'bol_opr': bol}
            except Exception as e:
                log(f"  OPR no-filter query error on port {port}: {e}")

        else:
            # ── Regular data-table path ────────────────────────────────────
            date_col = OPR_COL_DATE or 'Date'
            dax = (
                f"EVALUATE SELECTCOLUMNS("
                f"FILTER('{OPR_TABLE}', '{OPR_TABLE}'[{date_col}] = DATE({yr},{mo},{dy})),"
                f"\"bok_opr\", '{OPR_TABLE}'[{bok_measure}],"
                f"\"bol_opr\", '{OPR_TABLE}'[{bol_measure}])"
            )
            try:
                rows = run_dax(port, dax)
                if rows:
                    r = rows[0]
                    bok = to_pct(r.get('bok_opr'))
                    bol = to_pct(r.get('bol_opr'))
                    log(f"  OPR from '{OPR_TABLE}': BOK={bok}%  BOL={bol}%")
                    return {'bok_opr': bok, 'bol_opr': bol}
                else:
                    log(f"  OPR query on port {port}: no rows for {report_date}")
            except Exception as e:
                log(f"  OPR query error on port {port}: {e}")

    return None


def query_scrap(ports, report_date):
    """Query Scrap/Car ($) from Power BI Desktop.

    Auto-discovers the scrap table/measure if SCRAP_TABLE / SCRAP_MEASURE are
    not set — tries all common Volvo Body Shop naming patterns automatically.
    Returns {'scrap_car': float|None} or None if not found.
    """
    yr, mo, dy = report_date.year, report_date.month, report_date.day

    # Build candidate (table, measure) pairs
    # If user has configured explicit values, use those first
    if SCRAP_TABLE and SCRAP_MEASURE:
        candidates = [(SCRAP_TABLE, SCRAP_MEASURE)]
    else:
        # Auto-discover: try all common Volvo/Body Shop scrap measure patterns
        candidates = [
            # Most likely: same card format as OPR
            ('Scrap Card',          'Scrap/Car'),
            ('Scrap Card',          'Scrap Cost/Car'),
            ('Scrap Card',          'Scrap Cost per Car'),
            ('Scrap Card',          'Scrap BOK Card'),
            ('Scrap BOK Card',      'Scrap BOK Card'),
            ('Scrap BOK Card',      'Scrap/Car'),
            ('BOK Scrap Card',      'BOK Scrap Card'),
            # Possibly in the same table as OPR
            (OPR_TABLE,             'Scrap/Car'),
            (OPR_TABLE,             'Scrap Cost/Car'),
            (OPR_TABLE,             'Scrap Cost per Car'),
            # Generic names
            ('Scrap',               'Scrap/Car'),
            ('Scrap',               'Scrap Cost per Car'),
            ('Scrap per Car',       'Scrap per Car'),
            ('Scrap Cost',          'Scrap Cost per Car'),
            ('Body Shop Scrap',     'Scrap/Car'),
            ('BOK KPI',             'Scrap/Car'),
        ]
        # Remove duplicates while preserving order
        seen = set(); candidates = [c for c in candidates if not (c in seen or seen.add(c))]

    date_tables = [SCRAP_DATE_TABLE] if SCRAP_DATE_TABLE else [
        OPR_DATE_TABLE, 'Date Table', 'Date', 'Calendar', 'DimDate', 'Dates', 'DateTable'
    ]
    date_tables = [d for d in date_tables if d]

    def _parse_val(r):
        val = r.get('scrap_car') or r.get('[scrap_car]')
        if val is None: return None
        try:
            return round(float(val), 4)
        except (TypeError, ValueError):
            return None

    for port in ports:
        for tbl, msr in candidates:
            # Try with date filter first
            for dt in date_tables:
                dax = (
                    f"EVALUATE CALCULATETABLE("
                    f"ROW(\"scrap_car\", '{tbl}'[{msr}]),"
                    f"'{dt}'[Date] = DATE({yr},{mo},{dy}))"
                )
                try:
                    rows = run_dax(port, dax)
                    if rows:
                        fval = _parse_val(rows[0])
                        if fval is not None:
                            log(f"  Scrap/Car '{tbl}'[{msr}] (dt='{dt}'): ${fval}")
                            # If this was auto-discovered, log the config hint
                            if not (SCRAP_TABLE and SCRAP_MEASURE):
                                log(f"  HINT: set SCRAP_TABLE='{tbl}' SCRAP_MEASURE='{msr}' to skip auto-discovery")
                            return {'scrap_car': fval}
                except Exception as e:
                    err = str(e)
                    if any(kw in err for kw in ('cannot be found','not found','does not exist')):
                        break  # wrong table, skip all date tables for this combo
                    # Other error: continue to next date table

            # No-date-filter fallback for this table/measure
            dax_nf = f"EVALUATE ROW(\"scrap_car\", '{tbl}'[{msr}])"
            try:
                rows = run_dax(port, dax_nf)
                if rows:
                    fval = _parse_val(rows[0])
                    if fval is not None:
                        log(f"  Scrap/Car '{tbl}'[{msr}] (no date filter): ${fval}")
                        if not (SCRAP_TABLE and SCRAP_MEASURE):
                            log(f"  HINT: set SCRAP_TABLE='{tbl}' SCRAP_MEASURE='{msr}' to skip auto-discovery")
                        return {'scrap_car': fval}
            except Exception:
                pass

    return None


def query_powerbi(report_date):
    """
    Query Power BI Desktop for all KPI data for report_date.
    Tries ALL running PBI instances and uses the first one that has A Shop data.
    Returns dict with kpis + ppt item arrays, or None if PBI not available.
    """
    ports = find_pbi_ports()
    if not ports:
        return None

    yr, mo, dy = report_date.year, report_date.month, report_date.day

    # Probe each port: use the first one that returns A Shop rows for this date
    port = None
    probe_dax = f"""
        EVALUATE SELECTCOLUMNS(
          FILTER('A Shop Body Count-FTT/DPV',
                 'A Shop Body Count-FTT/DPV'[Date] = DATE({yr},{mo},{dy})),
          "Model", 'A Shop Body Count-FTT/DPV'[Model])"""
    for p in ports:
        probe = run_dax(p, probe_dax)
        if probe:
            port = p
            log(f"Using PBI port {port} — has {len(probe)} A Shop row(s) for {report_date}")
            break
        else:
            log(f"Port {p} — no A Shop data for {report_date}, trying next...")

    if port is None:
        # No port returned data for this specific date — fall back to first available port
        port = ports[0]
        log(f"No port had data for {report_date}; defaulting to port {port}")

    log(f"Querying Power BI Desktop (port {port}) for {report_date}...")

    # ── A Shop summary ────────────────────────────────────────────────────────
    ashop_rows = run_dax(port, f"""
        EVALUATE SELECTCOLUMNS(
          FILTER('A Shop Body Count-FTT/DPV',
                 'A Shop Body Count-FTT/DPV'[Date] = DATE({yr},{mo},{dy})),
          "Model",        'A Shop Body Count-FTT/DPV'[Model],
          "SentToPaint",  'A Shop Body Count-FTT/DPV'[Sent to Paint],
          "SentToRepair", 'A Shop Body Count-FTT/DPV'[Sent to Repair],
          "FTT",          'A Shop Body Count-FTT/DPV'[FTT],
          "DefectCount",  'A Shop Body Count-FTT/DPV'[Defect Count],
          "DPV",          'A Shop Body Count-FTT/DPV'[DPV],
          "WG_DPV",       'A Shop Body Count-FTT/DPV'[W&G DPV],
          "WG_Defects",   'A Shop Body Count-FTT/DPV'[W&G Defect Count],
          "HOP_DPV",      'A Shop Body Count-FTT/DPV'[HOP DPV],
          "HOP_Defects",  'A Shop Body Count-FTT/DPV'[HOP Defect Count]
        )
    """)
    log(f"  A Shop summary: {len(ashop_rows)} model rows")

    # ── A Shop FTT defect items ────────────────────────────────────────────────
    ashop_ftt_rows = run_dax(port, f"""
        EVALUATE SELECTCOLUMNS(
          FILTER('A Shop Defects',
                 'A Shop Defects'[Date] = DATE({yr},{mo},{dy})),
          "body",       'A Shop Defects'[Body Number],
          "rfid",       'A Shop Defects'[RFID],
          "desc",       'A Shop Defects'[Item Description],
          "model",      'A Shop Defects'[Model],
          "link_stn",   'A Shop Defects'[Linking Station],
          "link_time",  'A Shop Defects'[Link Time],
          "close_stn",  'A Shop Defects'[Closing Station],
          "close_time", 'A Shop Defects'[Close Time],
          "location",   'A Shop Defects'[Location],
          "extra",      'A Shop Defects'[Extra Info]
        )
    """)
    log(f"  A Shop FTT items: {len(ashop_ftt_rows)}")

    # Format datetimes
    for r in ashop_ftt_rows:
        r['link_time']  = fmt_dt(r.get('link_time'))
        r['close_time'] = fmt_dt(r.get('close_time'))

    # ── B Shop (WD6) summary ──────────────────────────────────────────────────
    # Date column in B Shop is a String — filter by string match
    date_str = report_date.strftime('%Y-%m-%d')
    bshop_rows = run_dax(port, f"""
        EVALUATE SELECTCOLUMNS(
          FILTER('B Shop Body Count-FTT/DPV',
                 LEFT('B Shop Body Count-FTT/DPV'[Date], 10) = "{date_str}"),
          "Model",      'B Shop Body Count-FTT/DPV'[Model],
          "BodyOK",     'B Shop Body Count-FTT/DPV'[Body OK],
          "BIWOut",     'B Shop Body Count-FTT/DPV'[WD6 BIW Out],
          "WD6_FTT",    'B Shop Body Count-FTT/DPV'[WD6 FTT],
          "Defects",    'B Shop Body Count-FTT/DPV'[Defect Count],
          "WD6_DPV",    'B Shop Body Count-FTT/DPV'[WD6 DPV]
        )
    """)
    log(f"  B Shop summary: {len(bshop_rows)} model rows")

    # ── WD6 FTT items ─────────────────────────────────────────────────────────
    wd6_ftt_rows = run_dax(port, f"""
        EVALUATE SELECTCOLUMNS(
          FILTER('WD6 FTT Items',
                 'WD6 FTT Items'[Date] = DATE({yr},{mo},{dy})),
          "body",       'WD6 FTT Items'[Body Number],
          "rfid",       'WD6 FTT Items'[RFID],
          "desc",       'WD6 FTT Items'[Item Description],
          "model",      'WD6 FTT Items'[Model],
          "link_stn",   'WD6 FTT Items'[Linking Station],
          "link_time",  'WD6 FTT Items'[Link Time],
          "close_stn",  'WD6 FTT Items'[Closing Station],
          "close_time", 'WD6 FTT Items'[Close Time],
          "location",   'WD6 FTT Items'[Location],
          "extra",      'WD6 FTT Items'[Extra Info]
        )
    """)
    log(f"  WD6 FTT items: {len(wd6_ftt_rows)}")
    for r in wd6_ftt_rows:
        r['link_time']  = fmt_dt(r.get('link_time'))
        r['close_time'] = fmt_dt(r.get('close_time'))

    # ── WD6 DPV items ─────────────────────────────────────────────────────────
    wd6_dpv_rows = run_dax(port, f"""
        EVALUATE SELECTCOLUMNS(
          FILTER('WD6 DPV Items',
                 'WD6 DPV Items'[Date] = DATE({yr},{mo},{dy})),
          "body",      'WD6 DPV Items'[Body Number],
          "rfid",      'WD6 DPV Items'[RFID],
          "desc",      'WD6 DPV Items'[Item Description],
          "model",     'WD6 DPV Items'[Model],
          "station",   'WD6 DPV Items'[Linking Station],
          "location",  'WD6 DPV Items'[Location],
          "extra",     'WD6 DPV Items'[Extra Info]
        )
    """)
    log(f"  WD6 DPV items: {len(wd6_dpv_rows)}")

    # ── A Shop W&G DPV items (Weld/Geometry defects) ──────────────────────────
    # Table name may vary — try known candidates, fail gracefully if not found
    wg_dpv_rows = []
    for wg_tbl in ['A Shop W&G Defects', 'A Shop WG Defects', 'WG Defects Linked', 'A Shop W&G Items']:
        try:
            wg_dpv_rows = run_dax(port, f"""
                EVALUATE SELECTCOLUMNS(
                  FILTER('{wg_tbl}',
                         '{wg_tbl}'[Date] = DATE({yr},{mo},{dy})),
                  "body",      '{wg_tbl}'[Body Number],
                  "rfid",      '{wg_tbl}'[RFID],
                  "desc",      '{wg_tbl}'[Item Description],
                  "model",     '{wg_tbl}'[Model],
                  "station",   '{wg_tbl}'[Station],
                  "location",  '{wg_tbl}'[Location]
                )
            """)
            log(f"  A Shop W&G DPV items ({wg_tbl}): {len(wg_dpv_rows)}")
            break  # found the right table
        except Exception as e:
            log(f"  W&G table '{wg_tbl}' not found: {e}")
    if not wg_dpv_rows:
        log("  W&G DPV items: 0 (no matching PBI table — defect count from summary only)")

    # ── C Shop summary ─────────────────────────────────────────────────────────
    cshop_rows = run_dax(port, f"""
        EVALUATE SELECTCOLUMNS(
          FILTER('C Shop Body Count-FTT',
                 'C Shop Body Count-FTT'[Date] = DATE({yr},{mo},{dy})),
          "Model",       'C Shop Body Count-FTT'[Model],
          "CAL_OK",      'C Shop Body Count-FTT'[CAL Body OK],
          "FN1_OK",      'C Shop Body Count-FTT'[Final 1 Body OK],
          "FN2_OK",      'C Shop Body Count-FTT'[Final 2 Body OK],
          "CAL_FTT",     'C Shop Body Count-FTT'[CAL FTT],
          "Final1_FTT",  'C Shop Body Count-FTT'[Final 1 FTT],
          "EOL_FTT",     'C Shop Body Count-FTT'[EOL FTT],
          "FTT_CAL",     'C Shop Body Count-FTT'[FTT CAL],
          "FTT_FN1",     'C Shop Body Count-FTT'[FTT FINAL1],
          "FTT_EOL",     'C Shop Body Count-FTT'[FTT EOL],
          "CAL_DPV",     'C Shop Body Count-FTT'[CAL DPV],
          "FN1_DPV",     'C Shop Body Count-FTT'[FN1 DPV],
          "FN2_DPV",     'C Shop Body Count-FTT'[FN2 DPV]
        )
    """)
    log(f"  C Shop summary: {len(cshop_rows)} model rows")

    # ── C Shop defect items (CAL + Final line items) ───────────────────────────
    cshop_defect_rows = run_dax(port, f"""
        EVALUATE SELECTCOLUMNS(
          FILTER('C Shop Defects Linked',
                 YEAR('C Shop Defects Linked'[Link Time]) = {yr} &&
                 MONTH('C Shop Defects Linked'[Link Time]) = {mo} &&
                 DAY('C Shop Defects Linked'[Link Time]) = {dy}),
          "body",       'C Shop Defects Linked'[Body Number],
          "rfid",       'C Shop Defects Linked'[RFID],
          "desc",       'C Shop Defects Linked'[Item Description],
          "model",      'C Shop Defects Linked'[Model],
          "link_stn",   'C Shop Defects Linked'[Linking Station],
          "link_time",  'C Shop Defects Linked'[Link Time],
          "close_stn",  'C Shop Defects Linked'[Closing Station],
          "close_time", 'C Shop Defects Linked'[Close Time],
          "location",   'C Shop Defects Linked'[Location],
          "extra",      'C Shop Defects Linked'[Extra Info Link]
        )
    """)
    log(f"  C Shop defect items: {len(cshop_defect_rows)}")
    for r in cshop_defect_rows:
        r['link_time']  = fmt_dt(r.get('link_time'))
        r['close_time'] = fmt_dt(r.get('close_time'))

    # ── BOK / BOL OPR (from BIW SQD Axxos — may be a different PBI instance) ──
    opr   = query_opr(ports, report_date)
    scrap = query_scrap(ports, report_date)

    return {
        'ashop':         ashop_rows,
        'ashop_ftt':     ashop_ftt_rows,
        'bshop':         bshop_rows,
        'wd6_ftt':       wd6_ftt_rows,
        'wd6_dpv':       wd6_dpv_rows,
        'wg_dpv':        wg_dpv_rows,
        'cshop':         cshop_rows,
        'cshop_defects': cshop_defect_rows,
        'bok_opr':       opr.get('bok_opr')   if opr   else None,
        'bol_opr':       opr.get('bol_opr')   if opr   else None,
        'scrap_car':     scrap.get('scrap_car') if scrap else None,
    }


def build_kpis_from_pbi(pbi):
    """Build the kpis dict for MM_DATA from Power BI query results."""
    kpis = {k: {'val': None} for k in TARGETS}

    # ── BOK / BOL OPR ─────────────────────────────────────────────────────────
    if pbi.get('bok_opr') is not None:
        kpis['bok_opr'] = {'val': pbi['bok_opr']}
    if pbi.get('bol_opr') is not None:
        kpis['bol_opr'] = {'val': pbi['bol_opr']}

    # ── Scrap / Car ───────────────────────────────────────────────────────────
    if pbi.get('scrap_car') is not None:
        kpis['scrap_car'] = {'val': pbi['scrap_car']}

    # ── A Shop ────────────────────────────────────────────────────────────────
    if pbi['ashop']:
        # Aggregate across models (EX90 + PSTR)
        total_paint  = sum(r.get('SentToPaint',0)  or 0 for r in pbi['ashop'])
        total_repair = sum(r.get('SentToRepair',0) or 0 for r in pbi['ashop'])
        total_eol    = total_paint + total_repair
        total_def    = sum(r.get('DefectCount',0)  or 0 for r in pbi['ashop'])
        wg_def       = sum(r.get('WG_Defects',0)   or 0 for r in pbi['ashop'])

        ftt_val = round(total_paint / total_eol * 100, 4) if total_eol > 0 else None
        dpv_val = round(total_def  / total_eol, 4)        if total_eol > 0 else None
        wg_dpv  = round(wg_def     / total_eol, 4)        if total_eol > 0 else None

        kpis['ashop_ftt'] = {'val': ftt_val, 'eol': total_eol, 'ok': total_paint, 'repair': total_repair}
        kpis['ashop_dpv'] = {'val': dpv_val, 'bodies': total_eol, 'defects': int(total_def)}
        kpis['wg_dpv']    = {'val': wg_dpv,  'bodies': total_eol, 'defects': int(wg_def)}

    # ── B Shop (WD6) ──────────────────────────────────────────────────────────
    if pbi['bshop']:
        total_ok    = sum(r.get('BodyOK',0)  or 0 for r in pbi['bshop'])
        total_biw   = sum(r.get('BIWOut',0)  or 0 for r in pbi['bshop'])
        total_def   = sum(r.get('Defects',0) or 0 for r in pbi['bshop'])
        total_repair_wd6 = total_biw - total_ok if total_biw > total_ok else 0
        ftt_wd6 = round(total_ok / total_biw * 100, 4) if total_biw > 0 else None
        dpv_wd6 = round(total_def / total_biw, 4)       if total_biw > 0 else None
        kpis['wd6_ftt'] = {'val': ftt_wd6, 'eol': total_biw, 'ok': total_ok, 'repair': total_repair_wd6}
        kpis['wd6_dpv'] = {'val': dpv_wd6, 'bodies': total_biw, 'defects': int(total_def)}

    # ── C Shop ────────────────────────────────────────────────────────────────
    if pbi['cshop']:
        # Aggregate CAL, Final 1, Final 2 across models
        cal_ok  = sum(r.get('CAL_OK',0)  or 0 for r in pbi['cshop'])
        fn1_ok  = sum(r.get('FN1_OK',0)  or 0 for r in pbi['cshop'])
        fn2_ok  = sum(r.get('FN2_OK',0)  or 0 for r in pbi['cshop'])
        cal_eol = sum(r.get('FTT_CAL',0) or 0 for r in pbi['cshop'])
        fn1_eol = sum(r.get('FTT_FN1',0) or 0 for r in pbi['cshop'])
        fn2_eol = sum(r.get('FTT_EOL',0) or 0 for r in pbi['cshop'])

        cal_ftt = round(cal_ok / cal_eol * 100, 4) if cal_eol > 0 else None
        fn1_ftt = round(fn1_ok / fn1_eol * 100, 4) if fn1_eol > 0 else None
        fn2_ftt = round(fn2_ok / fn2_eol * 100, 4) if fn2_eol > 0 else None

        # Final 2 EOL = sum of FTT_EOL (total bodies at EOL gate)
        kpis['cal_ftt']    = {'val': cal_ftt, 'eol': cal_eol, 'ok': cal_ok,  'repair': cal_eol - cal_ok}
        kpis['final1_ftt'] = {'val': fn1_ftt, 'eol': fn1_eol, 'ok': fn1_ok,  'repair': fn1_eol - fn1_ok}
        kpis['final2_ftt'] = {'val': fn2_ftt, 'eol': fn2_eol, 'ok': fn2_ok,  'repair': fn2_eol - fn2_ok}

    return kpis


def build_ppt_items_from_pbi(pbi):
    """Build the ppt item arrays for MM_DATA from Power BI query results."""
    def split_by_model(rows, model_col='model'):
        e536 = [r for r in rows if '536' in str(r.get('model','')) or 'EX90' in str(r.get('model',''))]
        e519 = [r for r in rows if '519' in str(r.get('model','')) or 'PSTR' in str(r.get('model',''))]
        # If no split by model, put all in 536
        if not e536 and not e519: e536 = rows
        return e536, e519

    # ── A-Shop: split 'A Shop Defects' rows into three purposeful lists ──────────
    # All items from 'A Shop Defects' for the report date
    ashop_all  = pbi.get('ashop_ftt', [])
    # FTT repair items: only bodies closed at an HVR repair station
    # (matches the PBI 536/519 FTT Items report filter)
    ashop_repair = [r for r in ashop_all
                    if 'HVR' in str(r.get('close_stn', '')).upper()]
    # DPV items: all items on finalized bodies (close_stn set, any station)
    # This approximates the DPV defect scope (repair bodies + bodies that passed)
    ashop_dpv_all = [r for r in ashop_all if r.get('close_stn') or r.get('close_time')]
    # W&G items: items linked at W&G stations (supplement dedicated W&G table query)
    ashop_wg = [r for r in ashop_all
                if any(x in str(r.get('link_stn', '')).upper() for x in ['W&G', 'WG'])]

    ashop_536,     ashop_519     = split_by_model(ashop_repair)
    ashop_dpv_536, ashop_dpv_519 = split_by_model(ashop_dpv_all)
    ashop_wg_536,  ashop_wg_519  = split_by_model(ashop_wg)

    wd6_ftt_536, wd6_ftt_519 = split_by_model(pbi['wd6_ftt'])
    wd6_dpv_536, wd6_dpv_519 = split_by_model(pbi['wd6_dpv'])

    # W&G DPV: merge dedicated W&G table items + W&G-station items from A-Shop Defects
    wg_ded_536, wg_ded_519   = split_by_model(pbi.get('wg_dpv', []))
    wg_dpv_536 = wg_ded_536 + [r for r in ashop_wg_536 if r not in wg_ded_536]
    wg_dpv_519 = wg_ded_519 + [r for r in ashop_wg_519 if r not in wg_ded_519]

    log(f"  A-Shop split: {len(ashop_repair)} repair items, {len(ashop_dpv_all)} DPV items, "
        f"{len(ashop_wg)} W&G-station items, {len(wg_dpv_536)+len(wg_dpv_519)} W&G total")

    # C Shop: split defects by linking station prefix to CAL vs Final
    cal_defects = [r for r in pbi['cshop_defects']
                   if any(x in str(r.get('link_stn','')).upper()
                          for x in ['CAL','EXCAL','EX-CAL','CALCAL'])]
    fn_defects  = [r for r in pbi['cshop_defects'] if r not in cal_defects]
    cal_536, cal_519   = split_by_model(cal_defects)
    fn2_536, fn2_519   = split_by_model(fn_defects)

    def _s(v, default=''):
        """Convert to string and strip newlines/carriage-returns (they break inline JS)."""
        return str(v if v is not None else default).replace('\r', '').replace('\n', ' ').strip()

    def to_ftt_item(r):
        return {
            'body':       _s(r.get('body'), '—'),
            'rfid':       _s(r.get('rfid'), '—'),
            'desc':       _s(r.get('desc')),
            'model':      _s(r.get('model')),
            'link_stn':   _s(r.get('link_stn')),
            'link_time':  _s(r.get('link_time')),
            'close_stn':  _s(r.get('close_stn')),
            'close_time': _s(r.get('close_time')),
            'location':   _s(r.get('location')),
            'extra':      _s(r.get('extra') or ''),
        }

    def to_dpv_item(r):
        return {
            'body':     _s(r.get('body'), '—'),
            'rfid':     _s(r.get('rfid'), '—'),
            'count':    1,
            'desc':     _s(r.get('desc')),
            'model':    _s(r.get('model')),
            'station':  _s(r.get('station', r.get('link_stn', ''))),
            'location': _s(r.get('location')),
            'extra':    _s(r.get('extra') or ''),
        }

    return {
        'ashop_ftt_536':  [to_ftt_item(r) for r in ashop_536],
        'ashop_ftt_519':  [to_ftt_item(r) for r in ashop_519],
        'ashop_dpv_536':  [to_ftt_item(r) for r in ashop_dpv_536],  # all-body defects for DPV table
        'ashop_dpv_519':  [to_ftt_item(r) for r in ashop_dpv_519],
        'wg_dpv_536':     [to_dpv_item(r) for r in wg_dpv_536],
        'wg_dpv_519':     [to_dpv_item(r) for r in wg_dpv_519],
        'wd6_ftt_536':   [to_ftt_item(r) for r in wd6_ftt_536],
        'wd6_ftt_519':   [to_ftt_item(r) for r in wd6_ftt_519],
        'wd6_dpv_536':   [to_dpv_item(r) for r in wd6_dpv_536],
        'wd6_dpv_519':   [to_dpv_item(r) for r in wd6_dpv_519],
        'cal_ftt_536':   [to_ftt_item(r) for r in cal_536],
        'cal_ftt_519':   [to_ftt_item(r) for r in cal_519],
        'final2_ftt_536':[to_ftt_item(r) for r in fn2_536],
        'final2_ftt_519':[to_ftt_item(r) for r in fn2_519],
    }


# ═══════════════════════════════════════════════════════════════════════════════
# DOWNTIME READING (Andon Downtime Sheet — same as Andon Dashboard)
# ═══════════════════════════════════════════════════════════════════════════════

def read_area_dt(hop_ws, dt_ws, wk_str, day_int):
    areas = {k: {'total':0,'planned':0,'unplanned':0,'events':[]}
             for k in ['underbody','upperbody','hangon','bodysides']}

    for ws, code_fn, area_map in [(hop_ws, hop_code, HOP_AREA),(dt_ws, dt_code, DT_AREA)]:
        if ws is None: continue
        for i, row in enumerate(ws.rows):
            if i == 0: continue
            vals = [c.value for c in row]
            try:    yr = int(float(str(vals[0])))
            except: continue
            if yr != 2026: continue
            if str(vals[1]).strip() != wk_str: continue
            try:    d = int(float(str(vals[2])))
            except: continue
            if d != day_int: continue

            start, end = vals[3], vals[4]
            dur = dur_min(start, end)
            if dur <= 0: continue

            area_str = str(vals[7]) if vals[7] and str(vals[7]) != 'None' else ''
            if not area_str: continue

            resp = vals[11] if len(vals) > 11 else None
            err  = vals[12] if len(vals) > 12 else None
            if should_exclude(resp, err): continue

            code     = code_fn(area_str)
            area_key = area_map.get(code)
            if not area_key: continue

            planned_flag = is_planned(resp, err)
            dur_r = round(dur, 1)
            areas[area_key]['total']     = round(areas[area_key]['total']     + dur_r, 1)
            if planned_flag: areas[area_key]['planned']  = round(areas[area_key]['planned']  + dur_r, 1)
            else:            areas[area_key]['unplanned']= round(areas[area_key]['unplanned']+ dur_r, 1)

            if len(areas[area_key]['events']) < 10:
                cause = ' · '.join(filter(None,[str(resp).strip() if resp else '',
                                                str(err).strip()  if err  else '']))[:60] or area_str[:30]
                areas[area_key]['events'].append({
                    'station': area_str[:20], 'start': t_fmt(start), 'end': t_fmt(end),
                    'dur_min': dur_r, 'cause': cause, 'planned': planned_flag
                })
    return areas


def build_hop_stops(hop_ws, dt_ws, wk_str, day_int):
    """Build top-10 HOP line stops sorted by duration (for BOK OPR display)."""
    stops = []
    area_label = {**{v: v.title() for v in HOP_AREA.values()},
                  **{v: v.title() for v in DT_AREA.values()}}
    area_label['underbody'] = 'Underbody'
    area_label['upperbody'] = 'Upperbody'
    area_label['hangon']    = 'Hang-On'
    area_label['bodysides'] = 'Body Sides'

    for ws, code_fn, area_map in [(hop_ws, hop_code, HOP_AREA),(dt_ws, dt_code, DT_AREA)]:
        if ws is None: continue
        for i, row in enumerate(ws.rows):
            if i == 0: continue
            vals = [c.value for c in row]
            try:    yr2 = int(float(str(vals[0])))
            except: continue
            if yr2 != 2026: continue
            if str(vals[1]).strip() != wk_str: continue
            try:    d = int(float(str(vals[2])))
            except: continue
            if d != day_int: continue

            start, end = vals[3], vals[4]
            dur = dur_min(start, end)
            if dur <= 0: continue

            area_str = str(vals[7]) if vals[7] else ''
            if not area_str: continue

            resp = vals[11] if len(vals) > 11 else None
            err  = vals[12] if len(vals) > 12 else None
            if should_exclude(resp, err): continue

            code     = code_fn(area_str)
            area_key = area_map.get(code, 'other')
            planned  = is_planned(resp, err)
            cause    = ' · '.join(filter(None,[str(resp).strip() if resp else '',
                                               str(err).strip()  if err  else '']))[:60] or area_str[:25]
            stops.append({
                'line':    area_str[:20],
                'area':    area_label.get(area_key, 'Other'),
                'dur_min': round(dur, 1),
                'cause':   cause,
                'planned': planned,
            })

    # Sort by duration DESC, deduplicate by line name keeping max duration
    seen = {}
    for s in sorted(stops, key=lambda x: -x['dur_min']):
        if s['line'] not in seen:
            seen[s['line']] = s
    return sorted(seen.values(), key=lambda x: -x['dur_min'])[:10]


# ═══════════════════════════════════════════════════════════════════════════════
# PPT READING (Safety, Part Quality, Bodies OOF, Scrap — table slides only)
# ═══════════════════════════════════════════════════════════════════════════════

def find_ppt(wk_str, day_int):
    if not os.path.isdir(PPT_DIR): return None
    wk_num = wk_str.replace('WK','')
    fname  = f"A Shop 26W{wk_num}D{day_int}.pptx"

    # Candidate folders: root, then week subfolders (e.g. 26WK12, WK12, 26W12)
    wk_subfolder_names = [f"26{wk_str}", wk_str, f"26W{wk_num}"]
    search_dirs = [PPT_DIR] + [
        os.path.join(PPT_DIR, sub) for sub in wk_subfolder_names
        if os.path.isdir(os.path.join(PPT_DIR, sub))
    ]

    for d in search_dirs:
        direct = os.path.join(d, fname)
        if os.path.exists(direct): return direct
        for f in os.listdir(d):
            if f.lower() == fname.lower():
                return os.path.join(d, f)

    log(f"WARNING: PPT not found: {fname}")
    return None

def ensure_pptx_deps():
    """Ensure python-pptx and markitdown[pptx] are installed."""
    try:
        from pptx import Presentation  # noqa
        return True
    except ImportError:
        log("python-pptx not installed — installing now...")
        try:
            subprocess.run(
                [sys.executable, '-m', 'pip', 'install', 'markitdown[pptx]', 'python-pptx'],
                capture_output=True, timeout=300)
            log("python-pptx installed OK")
            return True
        except Exception as e:
            log(f"WARNING: could not install python-pptx: {e}")
            return False

# Keep old name so any external callers still work
def ensure_markitdown():
    return ensure_pptx_deps()


def _ocr_slide_images(slide):
    """Try OCR on picture shapes in a slide. Returns list of text lines.
    Only used when a slide has NO text frames and NO tables (image-only slide).
    Requires pytesseract + Pillow; silently returns [] if not available."""
    try:
        import io as _io
        from PIL import Image as _Image
        import pytesseract as _tess
        text_lines = []
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                try:
                    img = _Image.open(_io.BytesIO(shape.image.blob))
                    # Downscale large images for faster OCR (max 1000px wide)
                    w, h = img.size
                    if w > 1000:
                        img = img.resize((1000, int(h * 1000 / w)), _Image.LANCZOS)
                    # Convert to RGB (some PNGs have RGBA transparency that slows OCR)
                    if img.mode not in ('L', 'RGB'):
                        img = img.convert('RGB')
                    ocr_text = _tess.image_to_string(img, config='--psm 3')
                    lines = [l.strip() for l in ocr_text.split('\n')
                             if l.strip() and len(l.strip()) > 3]
                    text_lines.extend(lines)
                except Exception:
                    pass
        return text_lines
    except ImportError:
        return []


def _inject_ocr_into_md(md_text, image_only_ocr):
    """Insert OCR lines for image-only slides into existing markitdown output."""
    import re as _re
    result = []
    chunks = _re.split(r'(<!--\s*Slide number:\s*\d+\s*-->)', md_text)
    for chunk in chunks:
        m = _re.match(r'<!--\s*Slide number:\s*(\d+)\s*-->', chunk.strip())
        result.append(chunk)
        if m:
            snum = int(m.group(1))
            if snum in image_only_ocr:
                result.append('\n' + '\n'.join(image_only_ocr[snum]) + '\n')
    return ''.join(result)


def read_ppt_markdown(ppt_path):
    """Read a .pptx file and return a markdown string with <!-- Slide number: N --> markers.

    Primary path:  markitdown Python API (markitdown[pptx])
    Fallback path: python-pptx direct read -> synthesises the same markdown format

    Either way: image-only slides (safety escalation reports, safety alerts)
    are supplemented with OCR text so the safety detector can find them.
    """
    ensure_pptx_deps()

    # Pre-scan for image-only slides and OCR them (works for both paths below)
    image_only_ocr = {}
    try:
        from pptx import Presentation as _Prs
        _prs = _Prs(ppt_path)
        for _idx, _slide in enumerate(_prs.slides):
            _snum = _idx + 1
            _has_text  = any(s.has_text_frame and s.text_frame.text.strip()
                             for s in _slide.shapes)
            _has_table = any(s.has_table for s in _slide.shapes)
            _has_pics  = any(s.shape_type == 13 for s in _slide.shapes)
            if _has_pics and not _has_text and not _has_table:
                _lines = _ocr_slide_images(_slide)
                if _lines:
                    image_only_ocr[_snum] = _lines
                    log(f"  PPT OCR slide {_snum}: {len(_lines)} lines extracted")
    except Exception as _e:
        log(f"  PPT OCR pre-scan: {_e}")

    # Primary: markitdown Python API
    try:
        from markitdown import MarkItDown
        md_converter = MarkItDown()
        result = md_converter.convert(ppt_path)
        md_text = result.text_content
        if md_text and md_text.strip():
            log(f"  PPT markitdown API: {len(md_text)} chars")
            if image_only_ocr:
                md_text = _inject_ocr_into_md(md_text, image_only_ocr)
            return md_text
    except Exception as e:
        log(f"  markitdown API unavailable ({type(e).__name__}) -- using python-pptx fallback")

    # Fallback: python-pptx direct -> mimic markitdown slide format
    try:
        from pptx import Presentation
        from pptx.util import Pt
        import io

        prs = Presentation(ppt_path)
        parts = []
        for idx, slide in enumerate(prs.slides):
            slide_num = idx + 1
            parts.append(f"\n<!-- Slide number: {slide_num} -->\n")

            for shape in slide.shapes:
                # Text frames
                if shape.has_text_frame:
                    txt = shape.text_frame.text.strip()
                    if txt:
                        try:
                            size = shape.text_frame.paragraphs[0].runs[0].font.size
                            is_heading = size and size >= Pt(18)
                        except Exception:
                            is_heading = False
                        if is_heading:
                            parts.append(f"## {txt}\n\n")
                        else:
                            parts.append(f"{txt}\n\n")

                # Tables -> markdown table
                if shape.has_table:
                    tbl = shape.table
                    rows_out = []
                    for row_idx, row in enumerate(tbl.rows):
                        cells = [c.text.strip().replace('\n', ' ') for c in row.cells]
                        rows_out.append('| ' + ' | '.join(cells) + ' |')
                        if row_idx == 0:
                            rows_out.append('| ' + ' | '.join(['---'] * len(cells)) + ' |')
                    parts.append('\n'.join(rows_out) + '\n\n')

            # Append OCR text for image-only slides
            if slide_num in image_only_ocr:
                parts.append('\n'.join(image_only_ocr[slide_num]) + '\n\n')

        md_text = ''.join(parts)
        log(f"  PPT python-pptx fallback: {len(md_text)} chars, {len(prs.slides)} slides")
        return md_text

    except Exception as e:
        log(f"WARNING: PPT read completely failed: {e}")
        return ''

def parse_md_table(text):
    """Parse a markdown pipe-table from text. Strips zero-width/invisible chars."""
    # Characters to strip beyond normal whitespace: zero-width space, ZWNJ, ZWJ, BOM
    _ZW = '\u200b\u200c\u200d\ufeff\u00ad'
    def _clean(s):
        return s.strip().strip(_ZW).strip()
    rows, headers, in_table = [], [], False
    for line in text.split('\n'):
        line = line.strip()
        if not (line.startswith('|') and line.endswith('|')):
            if in_table: break
            continue
        cells = [_clean(c) for c in line[1:-1].split('|')]
        if not in_table:
            headers = [_clean(h) for h in cells]; in_table = True; continue
        if all(re.match(r'^[-: ]+$', c) for c in cells if c): continue
        if any(c and c != 'None' for c in cells):
            rows.append({headers[j] if j < len(headers) else f'c{j}': cells[j]
                         for j in range(len(cells))})
    return rows

def _first_val(r, *keys):
    """Return first non-empty value from a dict matching any of the given keys (case-insensitive)."""
    ku = {k.upper(): v for k, v in r.items()}
    for key in keys:
        v = ku.get(key.upper(), '')
        if v: return v
    # Partial match fallback
    for key in keys:
        for k, v in ku.items():
            if key.upper() in k and v:
                return v
    return ''


def parse_ppt_tables(md_text):
    """Parse PPT markdown (all slides) using content-based detection.
    NEVER relies on hardcoded slide numbers — slides change position every day.
    Detects each section by table headers and text keywords."""
    data = {
        'safety':       {'title': '', 'detail': '', 'meta': ''},
        'part_quality': [],
        'bodies_oof':   [],
        'scrap':        '$0',
        'scrap_note':   'No scrap today',
    }
    if not md_text:
        return data

    chunks = re.split(r'<!--\s*Slide number:\s*(\d+)\s*-->', md_text)
    # Track which slides were already classified as data-table slides (PQ or BOF)
    # so we don't also detect them as safety/scrap
    data_table_slides = set()

    for i in range(1, len(chunks), 2):
        slide_num = int(chunks[i])
        content   = chunks[i + 1] if i + 1 < len(chunks) else ''
        cu        = content.upper()

        # ── Part Quality ────────────────────────────────────────────────────────
        # Detect by table having: PROD, PART DESCRIPTION, SUPPLIER, HOW BIG
        if not data['part_quality']:
            rows = parse_md_table(content)
            if rows:
                hdr_str = ' '.join(rows[0].keys()).upper()
                if ('PART DESCRIPTION' in hdr_str or
                        ('PROD' in hdr_str and 'SUPPLIER' in hdr_str)):
                    items = []
                    for r in rows:
                        part = _first_val(r, 'PART DESCRIPTION')
                        if not part:
                            continue
                        items.append({
                            'area':     _first_val(r, 'PROD.AREA', 'PROD. AREA', 'AREA', 'PRODUCTION AREA'),
                            'model':    _first_val(r, 'EX90 /723N', 'EX90/723N', 'MODEL', 'EX90'),
                            'part':     part,
                            'supplier': _first_val(r, 'SUPPLIER NAME', 'SUPPLIER'),
                            'partno':   _first_val(r, 'PART NUMBER', 'PART NO', 'PARTNO'),
                            'qty':      _first_val(r, 'HOW BIG/HOW MANY?', 'HOW BIG/HOW MANY', 'QTY', 'HOW BIG'),
                            'status':   _first_val(r, 'STATUS'),
                            'repeater': _first_val(r, 'REPEATER'),
                            'sort':     _first_val(r, 'SORT'),
                            'problem':  _first_val(r, 'PROBLEM STATEMENT/DETAILS', 'PROBLEM STATEMENT', 'PROBLEM'),
                            'handshake':_first_val(r, 'HANDSHAKE/VIRA', 'HANDSHAKE / VIRA', 'HANDSHAKE', 'VIRA'),
                        })
                    if items:
                        data['part_quality'] = items
                        data_table_slides.add(slide_num)
                        log(f"  PPT: Part Quality found on slide {slide_num} ({len(items)} items)")

        # ── Bodies Out of Flow ──────────────────────────────────────────────────
        # Detect by table having: RFID, LOCATION, AUTOMATION, BODY
        if not data['bodies_oof']:
            rows = parse_md_table(content)
            if rows:
                hdr_str = ' '.join(rows[0].keys()).upper()
                if ('RFID' in hdr_str or
                        ('LOCATION' in hdr_str and ('AUTOMATION' in hdr_str or 'BODY' in hdr_str))):
                    items = []
                    for r in rows:
                        rfid = _first_val(r, 'BODY/RFID', 'RFID', 'BODY')
                        items.append({
                            'mode':     _first_val(r, 'AUTOMATION /MANUAL', 'AUTOMATION/MANUAL', 'AUTOMATION', 'MANUAL'),
                            'rfid':     rfid,
                            'type':     _first_val(r, 'TYPE'),
                            'bodytype': _first_val(r, 'UNDERBODY TYPE', 'COMPLETE BODY', 'BODY TYPE') or 'Complete',
                            'location': _first_val(r, 'LOCATION STAGED', 'LOCATION'),
                            'status':   _first_val(r, 'STATUS'),
                            'reason':   _first_val(r, 'REASON'),
                            'removed':  _first_val(r, 'DATE REMOVED FROM LINE', 'DATE REMOVED', 'REMOVED'),
                            'expected': _first_val(r, 'EXPECTED REPAIR DATE', 'EXPECTED REPAIR', 'EXPECTED'),
                            'dummy':    _first_val(r, 'DUMMY ORDER ENTERED', 'DUMMY ORDER', 'DUMMY'),
                            'champion': _first_val(r, 'RESPONSIBLE CHAMPION', 'CHAMPION', 'RESPONSIBLE'),
                        })
                    if items:
                        data['bodies_oof'] = items
                        data_table_slides.add(slide_num)
                        log(f"  PPT: Bodies OOF found on slide {slide_num} ({len(items)} bodies)")

        # ── Safety ──────────────────────────────────────────────────────────────
        # Detect ONLY genuine safety cross / one-pager slides.
        # Must have an explicit safety alert phrase OR "SAFETY" + an injury word.
        # NEVER triggers on quality-only keywords like NOK, SUSPECT, GEMBA, ESCALAT.
        if not data['safety']['title'] and slide_num not in data_table_slides:
            # Only count keyword hits in NON-table lines
            non_table_lines = [l.strip() for l in content.split('\n')
                               if l.strip()
                               and not l.strip().startswith('|')   # table rows
                               and not l.strip().startswith('!')   # images
                               and not re.match(r'^[-|: ]+$', l.strip())]  # dividers
            non_table_text = ' '.join(non_table_lines).upper()
            # Strong phrase triggers (safety cross / one-pager / escalation report)
            SAFETY_STRONG = ['SAFETY ALERT', 'SAFETY CROSS', 'NEAR MISS',
                             'FIRST AID', 'RECORDABLE INCIDENT', 'STOP THE LINE',
                             'ESCALATION REPORT']  # Volvo safety escalation report template
            # Weak injury words — only trigger when "SAFETY" also present on slide
            SAFETY_INJURY = ['INJURY', 'INJURED', 'LACERATION', 'CUT/PUNCTURE',
                             'CUT/BURN', 'SLIP', 'FALL', 'STRUCK BY', 'PINCH POINT',
                             'ABRASION', 'STRAIN', 'CONTUSION']
            _is_safety_slide = (
                any(phrase in non_table_text for phrase in SAFETY_STRONG) or
                ('SAFETY' in non_table_text and
                 any(kw in non_table_text for kw in SAFETY_INJURY))
            )
            if _is_safety_slide:
                # PPT template placeholder patterns to skip
                TMPL_PATS = [
                    r'^hard subsection$', r'enter\s+pictures?\s+here', r'click\s+to\s+(add|edit)',
                    r'^add\s+your\s+text', r'^subtitle$', r'^type\s+here$',
                    r'^text\s+box$', r'^insert\s+picture', r'^picture\s+placeholder',
                    r'^\(enter\s+pictures?\s+here\)$',
                ]
                def _is_template(line):
                    l = line.strip().lower()
                    return any(re.search(p, l) for p in TMPL_PATS) or len(l) < 4
                # Get non-template text lines — strip any ## heading markers but keep content
                text_lines = []
                for _l in non_table_lines:
                    _stripped = re.sub(r'^#+\s*', '', _l).strip()
                    if _stripped and not _is_template(_stripped):
                        text_lines.append(_stripped)
                # Prefer lines that directly contain a safety phrase as the title
                title_line = ''
                detail_lines = []
                _ALL_SAFETY = SAFETY_STRONG + SAFETY_INJURY + ['SAFETY', 'ALERT']
                # Look for a line with a strong safety keyword first
                for l in text_lines:
                    if any(kw in l.upper() for kw in _ALL_SAFETY) and len(l) > 6:
                        title_line = l
                        break
                if not title_line and text_lines:
                    title_line = text_lines[0]
                for l in text_lines:
                    if l == title_line:
                        continue
                    if len(l) > 8:
                        detail_lines.append(l)
                    if len(detail_lines) >= 3:
                        break
                if title_line:
                    data['safety'] = {
                        'title':  title_line,
                        'detail': ' | '.join(detail_lines),
                        'meta':   f'Slide {slide_num}',
                    }
                    log(f"  PPT: Safety found on slide {slide_num}: {title_line[:60]}")

        # ── Scrap ───────────────────────────────────────────────────────────────
        # Detect by text containing SCRAP keyword (skip data-table slides)
        if data['scrap'] == '$0' and 'SCRAP' in cu and slide_num not in data_table_slides:
            # Match "$1,234" or "1,234$" or "0$" formats
            dollar = re.search(r'\$\s*[\d,]+', content) or re.search(r'[\d,]+\s*\$', content)
            if dollar:
                raw = dollar.group(0).replace(' ', '').strip('$').replace(',','')
                try:
                    amt = int(raw)
                except Exception:
                    amt = -1
                if amt == 0:
                    data['scrap'] = '$0'
                    data['scrap_note'] = 'No scrap today'
                    log(f"  PPT: Scrap slide {slide_num}: $0 (no scrap)")
                else:
                    scrap_val = '$' + dollar.group(0).replace(' ', '').strip('$')
                    data['scrap'] = scrap_val
                    data['scrap_note'] = f"{scrap_val} scrap cost — review breakdown with team"
                    log(f"  PPT: Scrap found on slide {slide_num}: {scrap_val}")
            elif 'NO SCRAP' in cu:
                data['scrap'] = '$0'
                data['scrap_note'] = 'No scrap today'
                log(f"  PPT: Scrap slide {slide_num}: NO SCRAP")

    # Log what was / wasn't found
    if not data['part_quality']:  log("  PPT WARNING: Part Quality table not found in any slide")
    if not data['bodies_oof']:    log("  PPT WARNING: Bodies OOF table not found in any slide")
    if not data['safety']['title']: log("  PPT WARNING: Safety action item not found in any slide")

    return data


# ═══════════════════════════════════════════════════════════════════════════════
# MM_DATA BUILDER
# ═══════════════════════════════════════════════════════════════════════════════

def build_mm_data(areas, hop_stops, kpis, ppt_items, ppt_tables, report_date, wk_str, day_int):
    SHIFT = 540
    total     = round(sum(a['total']     for a in areas.values()), 1)
    planned   = round(sum(a['planned']   for a in areas.values()), 1)
    unplanned = round(sum(a['unplanned'] for a in areas.values()), 1)
    avail     = round((SHIFT - unplanned) / SHIFT * 100, 1)

    def area_js(key):
        a = areas[key]
        evts = json.dumps(a['events'], default=str)
        return f'{{"total":{a["total"]},"planned":{a["planned"]},"unplanned":{a["unplanned"]},"events":{evts}}}'

    hop_js  = json.dumps(hop_stops,         default=str, ensure_ascii=False)
    ub_js   = json.dumps([s for s in hop_stops if s['area']=='Underbody'],
                         default=str, ensure_ascii=False)
    up_js   = json.dumps([s for s in hop_stops if s['area']=='Upperbody'],
                         default=str, ensure_ascii=False)
    kpis_js = json.dumps(kpis,  default=str, ensure_ascii=False)
    saf_js  = json.dumps(ppt_tables['safety'],       ensure_ascii=False)
    pq_js   = json.dumps(ppt_tables['part_quality'], ensure_ascii=False)
    bof_js  = json.dumps(ppt_tables['bodies_oof'],   ensure_ascii=False)

    # Merge ppt_items into ppt_tables for FTT/DPV
    ppt_full = {**ppt_tables, **ppt_items}
    ppt_full_js = json.dumps(ppt_full, default=str, ensure_ascii=False)

    # Also build a plain dict (used for MM_HISTORY)
    downtime_dict = {
        'total_min': total, 'planned_min': planned,
        'unplanned_min': unplanned, 'availability': avail,
        'areas': {k: dict(areas[k]) for k in areas},
        'hop_stops': hop_stops,
        'ub_stops': [s for s in hop_stops if s['area']=='Underbody'],
        'up_stops': [s for s in hop_stops if s['area']=='Upperbody'],
    }
    day_dict = {
        'date': str(report_date), 'wk': wk_str, 'day': day_int,
        'kpis': kpis, 'downtime': downtime_dict, 'ppt': ppt_full
    }

    js_str = (
        f"const MM_DATA = {{\n"
        f"  date: {json.dumps(str(report_date))},\n"
        f"  wk: {json.dumps(wk_str)},\n"
        f"  day: {day_int},\n"
        f"  kpis: {kpis_js},\n"
        f"  downtime: {{\n"
        f"    total_min: {total}, planned_min: {planned},\n"
        f"    unplanned_min: {unplanned}, availability: {avail},\n"
        f"    areas: {{\n"
        f"      underbody: {area_js('underbody')},\n"
        f"      upperbody: {area_js('upperbody')},\n"
        f"      hangon:    {area_js('hangon')},\n"
        f"      bodysides: {area_js('bodysides')}\n"
        f"    }},\n"
        f"    hop_stops: {hop_js},\n"
        f"    ub_stops:  {ub_js},\n"
        f"    up_stops:  {up_js}\n"
        f"  }},\n"
        f"  ppt: {ppt_full_js}\n"
        f"}};"
    )
    return js_str, day_dict


def patch_html(new_mm_js, day_dict, report_date, target_file=None):
    HISTORY_MAX = 10  # keep last 10 working days

    target = target_file or DASH
    with open(target, 'r', encoding='utf-8') as f: html = f.read()

    # ── patch MM_DATA ─────────────────────────────────────────────────────────
    new_html, n = re.subn(r'const MM_DATA = \{.*?\};', new_mm_js, html, flags=re.DOTALL)
    if n != 1:
        log(f"ERROR: Expected 1 MM_DATA block, got {n}")
        return False

    # ── patch MM_HISTORY ──────────────────────────────────────────────────────
    hist_match = re.search(r'const MM_HISTORY = \{.*?\};', new_html, flags=re.DOTALL)
    if hist_match:
        try:
            existing_js = hist_match.group(0)
            # Extract JSON body from the JS object literal
            body = re.sub(r'^const MM_HISTORY = ', '', existing_js).rstrip(';')
            history = json.loads(body)
        except Exception:
            history = {}
    else:
        history = {}

    # Add/update today's entry and trim to last HISTORY_MAX days
    history[str(report_date)] = day_dict
    sorted_dates = sorted(history.keys())
    if len(sorted_dates) > HISTORY_MAX:
        for old in sorted_dates[:-HISTORY_MAX]:
            del history[old]

    new_history_js = 'const MM_HISTORY = ' + json.dumps(history, default=str, ensure_ascii=False) + ';'

    if hist_match:
        new_html = new_html[:hist_match.start()] + new_history_js + new_html[hist_match.end():]
    else:
        # Fallback: append before AREA_CFG
        new_html = new_html.replace(
            '/* ─── Area display config ─── */',
            new_history_js + '\n\n/* ─── Area display config ─── */'
        )

    with open(target, 'w', encoding='utf-8') as f: f.write(new_html)
    log(f"HTML patched: {target}  (history: {len(history)} days)")
    return True


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def update():
    today       = datetime.date.today()
    report_date = prev_working_day(today)
    wk_str, day_int = date_to_wk_day(report_date)

    log(f"=== Morning Meeting update started ({today}) ===")
    log(f"Report date: {report_date}  ({wk_str} D{day_int})")
    if wk_str is None: log("ERROR: weekend"); return False

    # ── 1. Power BI query (FTT/DPV — pre-calculated, no math needed) ──────────
    pbi = query_powerbi(report_date)
    if pbi:
        kpis      = build_kpis_from_pbi(pbi)
        ppt_items = build_ppt_items_from_pbi(pbi)
        log(f"  Power BI data loaded successfully")
        for k, v in kpis.items():
            if isinstance(v, dict) and v.get('val') is not None:
                log(f"    {k}: {v['val']}")
    else:
        log("  WARNING: Power BI not available — using cached/sample KPI values")
        log("  TIP: Open Power BI Desktop, refresh the FTT-DPV file, then re-run this script")
        kpis      = {k: {'val': None} for k in TARGETS}
        ppt_items = {}

    # ── 2. Downtime from OneDrive Excel ───────────────────────────────────────
    tmp_hop = os.path.join(WORK, '_mm_hop.xlsm')
    tmp_dt  = os.path.join(WORK, '_mm_dt.xlsm')
    for src, dst, name in [(HOP_SRC, tmp_hop,'HOP'),(DT_SRC, tmp_dt,'DT')]:
        if os.path.exists(src):
            shutil.copy2(src, dst)
            log(f"Copied {name}: {os.path.getsize(dst)//1024} KB")
        else:
            log(f"WARNING: {name} not found: {src}")

    hop_ws = dt_ws = None
    try:
        wb = openpyxl.load_workbook(tmp_hop, data_only=True, read_only=True)
        hop_ws = wb['New(DT LOG)']; log("Opened HOP workbook")
    except Exception as e: log(f"HOP error: {e}")
    try:
        wb = openpyxl.load_workbook(tmp_dt, data_only=True, read_only=True)
        dt_ws = wb['New(DT LOG)']; log("Opened DT workbook")
    except Exception as e: log(f"DT error: {e}")

    areas     = read_area_dt(hop_ws, dt_ws, wk_str, day_int)
    hop_stops = build_hop_stops(hop_ws, dt_ws, wk_str, day_int)

    for k, a in areas.items():
        log(f"  {k:12s}: total={a['total']} min  unplanned={a['unplanned']}")
    log(f"  Top HOP stop: {hop_stops[0]['line'] if hop_stops else 'none'}")

    # ── 3. PPT tables (Safety, Part Quality, Bodies OOF, Scrap) ──────────────
    ppt_path = find_ppt(wk_str, day_int)
    if ppt_path:
        log(f"Found PPT: {ppt_path}")
        md_text    = read_ppt_markdown(ppt_path)
        ppt_tables = parse_ppt_tables(md_text)
    else:
        ppt_tables = {'safety':{'title':'','detail':'','meta':''},
                      'part_quality':[], 'bodies_oof':[], 'scrap':'$0',
                      'scrap_note':'PPT not found — update manually'}

    # ── 4. Build + patch ──────────────────────────────────────────────────────
    mm_js, day_dict = build_mm_data(areas, hop_stops, kpis, ppt_items, ppt_tables,
                                    report_date, wk_str, day_int)
    success = patch_html(mm_js, day_dict, report_date)
    patch_html(mm_js, day_dict, report_date, target_file=DASH_MOBILE)
    log("Mobile dashboard also patched.")

    # ── 5. Push both dashboards to GitHub Pages ───────────────────────────────
    if GITHUB_ENABLED:
        try:
            now_str = datetime.datetime.now().strftime('%H:%M')
            for lock_name in ['index.lock', 'HEAD.lock', 'config.lock', 'packed-refs.lock']:
                lock_path = os.path.join(GITHUB_REPO, '.git', lock_name)
                if os.path.exists(lock_path):
                    try: os.remove(lock_path); log(f"Removed stale git lock: {lock_name}")
                    except Exception as le: log(f"Could not remove lock {lock_name}: {le}")

            def run_git(args, allow_fail=False):
                cmd = ['git', '-C', GITHUB_REPO] + args
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
                out = (result.stdout + result.stderr).strip()
                if result.returncode == 0: log(f"Git OK: {' '.join(args)}")
                else: log(f"Git {'WARN' if allow_fail else 'FAIL'}: {' '.join(args)} -> {out[:200]}")
                return result.returncode == 0

            run_git(['fetch', 'origin', 'main'])
            run_git(['reset', '--mixed', 'origin/main'])
            run_git(['add', 'morning_meeting_dashboard.html', 'morning_meeting_mobile.html'])
            run_git(['commit', '-m', f'MM-update-{report_date}-{now_str}'], allow_fail=True)
            ok = run_git(['push', 'origin', 'main'])
            if not ok:
                log("Git push failed — retrying with force-with-lease")
                run_git(['fetch', 'origin', 'main'])
                run_git(['reset', '--mixed', 'origin/main'])
                run_git(['add', 'morning_meeting_dashboard.html', 'morning_meeting_mobile.html'])
                run_git(['commit', '-m', f'MM-update-{report_date}-retry'], allow_fail=True)
                run_git(['push', '--force-with-lease', 'origin', 'main'])
        except Exception as e:
            log(f"Git ERROR: {e}")
    else:
        log("GitHub push skipped (GITHUB_ENABLED=False).")

    for tmp in [tmp_hop, tmp_dt]:
        try: os.remove(tmp)
        except: pass

    log("=== Morning Meeting update complete ===\n")
    return success


def patch_history_only(day_dict, report_date, target_file=None):
    """Add a day_dict to MM_HISTORY without touching MM_DATA.
    Used for backfilling past days."""
    HISTORY_MAX = 10
    target = target_file or DASH
    with open(target, 'r', encoding='utf-8') as f: html = f.read()

    hist_match = re.search(r'const MM_HISTORY = \{.*?\};', html, flags=re.DOTALL)
    if hist_match:
        try:
            body = re.sub(r'^const MM_HISTORY = ', '', hist_match.group(0)).rstrip(';')
            history = json.loads(body)
        except Exception:
            history = {}
    else:
        history = {}

    # ── Merge: preserve existing real KPI values if new backfill got null data ──
    existing = history.get(str(report_date), {})
    if existing:
        def _has_val(kpis):
            return any((v or {}).get('val') is not None for v in (kpis or {}).values()
                       if isinstance(v, dict))
        # If existing entry has real KPIs but new backfill doesn't, keep the old KPIs
        if _has_val(existing.get('kpis')) and not _has_val(day_dict.get('kpis')):
            log(f"  Keeping existing KPIs for {report_date} (PBI returned 0 rows this run)")
            day_dict['kpis'] = existing['kpis']
        # If existing entry has real FTT/DPV item lists, keep them
        if existing.get('ppt'):
            for item_key in ['ashop_ftt_536','ashop_ftt_519','wd6_ftt_536','wd6_ftt_519',
                             'wd6_dpv_536','wd6_dpv_519','cal_ftt_536','cal_ftt_519',
                             'final2_ftt_536','final2_ftt_519','wg_items_536','wg_items_519']:
                if existing['ppt'].get(item_key) and not (day_dict.get('ppt') or {}).get(item_key):
                    if 'ppt' not in day_dict: day_dict['ppt'] = {}
                    day_dict['ppt'][item_key] = existing['ppt'][item_key]
        # Preserve downtime events if new run got zeros
        if (existing.get('downtime',{}).get('total_min',0) > 0 and
                day_dict.get('downtime',{}).get('total_min',0) == 0):
            day_dict['downtime'] = existing['downtime']
    history[str(report_date)] = day_dict
    sorted_dates = sorted(history.keys())
    if len(sorted_dates) > HISTORY_MAX:
        for old in sorted_dates[:-HISTORY_MAX]:
            del history[old]

    new_history_js = 'const MM_HISTORY = ' + json.dumps(history, default=str, ensure_ascii=False) + ';'

    if hist_match:
        new_html = html[:hist_match.start()] + new_history_js + html[hist_match.end():]
    else:
        new_html = html.replace(
            '/* ─── Area display config ─── */',
            new_history_js + '\n\n/* ─── Area display config ─── */'
        )

    with open(target, 'w', encoding='utf-8') as f: f.write(new_html)
    log(f"History-only patch: {target}  (now {len(history)} days)")
    return True


def backfill(target_date):
    """Backfill MM_HISTORY for a specific past date using Power BI + OneDrive.
    Does NOT change MM_DATA. Pushes to GitHub when all dates are done."""
    wk_str, day_int = date_to_wk_day(target_date)
    if wk_str is None:
        log(f"Skipping {target_date} — not a working day")
        return False

    log(f"=== Backfill history: {target_date}  ({wk_str} D{day_int}) ===")

    # 1. Power BI (same as daily update)
    pbi = query_powerbi(target_date)
    if pbi:
        kpis      = build_kpis_from_pbi(pbi)
        ppt_items = build_ppt_items_from_pbi(pbi)
        log("  Power BI data loaded")
        for k, v in kpis.items():
            if isinstance(v, dict) and v.get('val') is not None:
                log(f"    {k}: {v['val']}")
    else:
        log("  WARNING: Power BI not available — KPIs will be null for this day")
        kpis      = {k: {'val': None} for k in TARGETS}
        ppt_items = {}

    # 2. Downtime (same as daily update)
    tmp_hop = os.path.join(WORK, '_mm_hop.xlsm')
    tmp_dt  = os.path.join(WORK, '_mm_dt.xlsm')
    for src, dst, name in [(HOP_SRC, tmp_hop, 'HOP'), (DT_SRC, tmp_dt, 'DT')]:
        if os.path.exists(src):
            shutil.copy2(src, dst)
            log(f"  Copied {name}")
        else:
            log(f"  WARNING: {name} not found: {src}")

    hop_ws = dt_ws = None
    try:
        wb = openpyxl.load_workbook(tmp_hop, data_only=True, read_only=True)
        hop_ws = wb['New(DT LOG)']
    except Exception as e: log(f"  HOP error: {e}")
    try:
        wb = openpyxl.load_workbook(tmp_dt, data_only=True, read_only=True)
        dt_ws = wb['New(DT LOG)']
    except Exception as e: log(f"  DT error: {e}")

    areas     = read_area_dt(hop_ws, dt_ws, wk_str, day_int)
    hop_stops = build_hop_stops(hop_ws, dt_ws, wk_str, day_int)

    # 3. PPT (same as daily update)
    ppt_path = find_ppt(wk_str, day_int)
    if ppt_path:
        log(f"  Found PPT: {ppt_path}")
        ppt_tables = parse_ppt_tables(read_ppt_markdown(ppt_path))
    else:
        ppt_tables = {'safety': {'title': '', 'detail': '', 'meta': ''},
                      'part_quality': [], 'bodies_oof': [],
                      'scrap': '$0', 'scrap_note': 'PPT not found for this day'}

    # 4. Build day_dict and patch ONLY history (MM_DATA unchanged)
    _, day_dict = build_mm_data(areas, hop_stops, kpis, ppt_items, ppt_tables,
                                target_date, wk_str, day_int)

    patch_history_only(day_dict, target_date, DASH)
    patch_history_only(day_dict, target_date, DASH_MOBILE)

    for tmp in [tmp_hop, tmp_dt]:
        try: os.remove(tmp)
        except: pass

    log(f"=== Backfill complete: {target_date} ===\n")
    return True


if __name__ == '__main__':
    import argparse
    parser = argparse.ArgumentParser(description='A-Shop Morning Meeting Dashboard updater')
    parser.add_argument(
        '--date', metavar='YYYY-MM-DD',
        help='Backfill MM_HISTORY for a specific past date without changing MM_DATA'
    )
    parser.add_argument(
        '--backfill-week', action='store_true',
        help='Backfill all working days of the current WK12 (Mon–Thu) into MM_HISTORY'
    )
    parser.add_argument(
        '--discover-tables', action='store_true',
        help='List all PBI table names, columns, and measures (for OPR/BOK/BOL tables) '
             'from all running Power BI Desktop instances.  Run with BIW SQD Axxos open, '
             'then fill in OPR_TABLE / OPR_BOK_MEASURE / OPR_BOL_MEASURE at the top of this file.'
    )
    args = parser.parse_args()

    if args.discover_tables:
        discover_tables()
        sys.exit(0)

    elif args.date:
        target = datetime.date.fromisoformat(args.date)
        ok = backfill(target)
        # Push to GitHub after backfill
        if ok and GITHUB_ENABLED:
            try:
                subprocess.run(['git', '-C', GITHUB_REPO, 'add',
                                'morning_meeting_dashboard.html',
                                'morning_meeting_mobile.html'], check=True)
                subprocess.run(['git', '-C', GITHUB_REPO, 'commit', '-m',
                                f'backfill-history-{args.date}'], check=True)
                subprocess.run(['git', '-C', GITHUB_REPO, 'push', 'origin', 'main'], check=True)
                log("GitHub push OK")
            except Exception as e:
                log(f"Git push error: {e}")
        sys.exit(0 if ok else 1)

    elif args.backfill_week:
        today  = datetime.date.today()
        mon    = today - datetime.timedelta(days=today.weekday())   # Monday of current week
        days   = [mon + datetime.timedelta(days=i) for i in range(5)
                  if (mon + datetime.timedelta(days=i)) < today]    # Mon–yesterday
        log(f"Backfilling {len(days)} day(s): {[str(d) for d in days]}")
        for d in days:
            backfill(d)
        # Single GitHub push after all days
        if GITHUB_ENABLED:
            try:
                subprocess.run(['git', '-C', GITHUB_REPO, 'add',
                                'morning_meeting_dashboard.html',
                                'morning_meeting_mobile.html'], check=True)
                subprocess.run(['git', '-C', GITHUB_REPO, 'commit', '-m',
                                f'backfill-history-week-{mon}'], check=True)
                subprocess.run(['git', '-C', GITHUB_REPO, 'push', 'origin', 'main'], check=True)
                log("GitHub push OK")
            except Exception as e:
                log(f"Git push error: {e}")
        sys.exit(0)

    else:
        sys.exit(0 if update() else 1)
